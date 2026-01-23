"""
Office Document Processor

Handles DOCX, XLSX, and PPTX file processing.
"""

import io
import re
import logging
from pathlib import Path
from typing import Any, Dict, List, Optional, Union

logger = logging.getLogger("docs_mcp.processors.office")


class OfficeProcessor:
    """
    Processor for Microsoft Office documents.
    """

    # =========================================================================
    # DOCX Processing
    # =========================================================================

    def read_docx(self, file_path: Union[str, Path, bytes]) -> Dict[str, Any]:
        """
        Extract text, tables, and structure from a DOCX file.

        Args:
            file_path: Path to the DOCX file or bytes content.

        Returns:
            Dict containing text, tables, and metadata.
        """
        try:
            from docx import Document
            from docx.table import Table
        except ImportError:
            return {"error": "python-docx not installed. Run: pip install python-docx"}

        try:
            if isinstance(file_path, bytes):
                doc = Document(io.BytesIO(file_path))
            else:
                doc = Document(file_path)

            # Extract text from paragraphs
            paragraphs = []
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append({
                        "text": para.text,
                        "style": para.style.name if para.style else None,
                    })

            # Extract tables
            tables = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                tables.append(table_data)

            # Extract metadata
            core_props = doc.core_properties
            metadata = {
                "author": core_props.author,
                "title": core_props.title,
                "subject": core_props.subject,
                "created": str(core_props.created) if core_props.created else None,
                "modified": str(core_props.modified) if core_props.modified else None,
            }

            # Full text extraction
            full_text = "\n".join(p["text"] for p in paragraphs)

            return {
                "text": full_text,
                "paragraphs": paragraphs,
                "tables": tables,
                "metadata": metadata,
                "paragraph_count": len(paragraphs),
                "table_count": len(tables),
            }

        except Exception as e:
            logger.error(f"Error reading DOCX: {e}")
            return {"error": str(e)}

    def fill_docx_template(
        self,
        template_path: Union[str, Path, bytes],
        data: Dict[str, str],
        output_path: Optional[Union[str, Path]] = None,
    ) -> Dict[str, Any]:
        """
        Replace {{placeholders}} in a DOCX with data values.

        Args:
            template_path: Path to the template DOCX or bytes content.
            data: Dict of placeholder names to values.
            output_path: Where to save the result (optional, returns bytes if not provided).

        Returns:
            Dict with success status and output path or content.
        """
        try:
            from docx import Document
        except ImportError:
            return {"error": "python-docx not installed. Run: pip install python-docx"}

        try:
            if isinstance(template_path, bytes):
                doc = Document(io.BytesIO(template_path))
            else:
                doc = Document(template_path)

            # Replace in paragraphs
            for para in doc.paragraphs:
                for key, value in data.items():
                    placeholder = "{{" + key + "}}"
                    if placeholder in para.text:
                        # Replace while preserving formatting
                        for run in para.runs:
                            if placeholder in run.text:
                                run.text = run.text.replace(placeholder, str(value))

            # Replace in tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for para in cell.paragraphs:
                            for key, value in data.items():
                                placeholder = "{{" + key + "}}"
                                if placeholder in para.text:
                                    for run in para.runs:
                                        if placeholder in run.text:
                                            run.text = run.text.replace(placeholder, str(value))

            # Replace in headers and footers
            for section in doc.sections:
                for header in [section.header, section.first_page_header, section.even_page_header]:
                    if header:
                        for para in header.paragraphs:
                            for key, value in data.items():
                                placeholder = "{{" + key + "}}"
                                if placeholder in para.text:
                                    for run in para.runs:
                                        if placeholder in run.text:
                                            run.text = run.text.replace(placeholder, str(value))

                for footer in [section.footer, section.first_page_footer, section.even_page_footer]:
                    if footer:
                        for para in footer.paragraphs:
                            for key, value in data.items():
                                placeholder = "{{" + key + "}}"
                                if placeholder in para.text:
                                    for run in para.runs:
                                        if placeholder in run.text:
                                            run.text = run.text.replace(placeholder, str(value))

            if output_path:
                doc.save(output_path)
                return {"success": True, "output_path": str(output_path)}
            else:
                buffer = io.BytesIO()
                doc.save(buffer)
                buffer.seek(0)
                return {"success": True, "content": buffer.read()}

        except Exception as e:
            logger.error(f"Error filling DOCX template: {e}")
            return {"error": str(e)}

    def docx_to_markdown(self, file_path: Union[str, Path, bytes]) -> Dict[str, Any]:
        """
        Convert DOCX to Markdown format.

        Args:
            file_path: Path to the DOCX file or bytes content.

        Returns:
            Dict containing the markdown text.
        """
        result = self.read_docx(file_path)
        if "error" in result:
            return result

        try:
            lines = []

            for para in result.get("paragraphs", []):
                text = para["text"]
                style = para.get("style", "")

                # Convert headings
                if style and "Heading" in style:
                    try:
                        level = int(style.replace("Heading ", "").replace("Heading", "1"))
                        lines.append("#" * level + " " + text)
                    except ValueError:
                        lines.append(text)
                elif style == "Title":
                    lines.append("# " + text)
                elif style == "List Paragraph":
                    lines.append("- " + text)
                else:
                    lines.append(text)

                lines.append("")  # Blank line between paragraphs

            # Add tables as markdown tables
            for i, table in enumerate(result.get("tables", [])):
                if table:
                    lines.append("")
                    # Header row
                    lines.append("| " + " | ".join(table[0]) + " |")
                    lines.append("| " + " | ".join(["---"] * len(table[0])) + " |")
                    # Data rows
                    for row in table[1:]:
                        lines.append("| " + " | ".join(row) + " |")
                    lines.append("")

            return {
                "markdown": "\n".join(lines),
                "metadata": result.get("metadata", {}),
            }

        except Exception as e:
            logger.error(f"Error converting DOCX to markdown: {e}")
            return {"error": str(e)}

    # =========================================================================
    # XLSX Processing
    # =========================================================================

    def read_xlsx(
        self,
        file_path: Union[str, Path, bytes],
        sheet_name: Optional[str] = None,
    ) -> Dict[str, Any]:
        """
        Read spreadsheet data from an XLSX file.

        Args:
            file_path: Path to the XLSX file or bytes content.
            sheet_name: Specific sheet to read (optional, reads all if not specified).

        Returns:
            Dict containing sheets data and metadata.
        """
        try:
            from openpyxl import load_workbook
        except ImportError:
            return {"error": "openpyxl not installed. Run: pip install openpyxl"}

        try:
            if isinstance(file_path, bytes):
                wb = load_workbook(io.BytesIO(file_path), data_only=True)
            else:
                wb = load_workbook(file_path, data_only=True)

            sheets = {}
            sheet_names = [sheet_name] if sheet_name else wb.sheetnames

            for name in sheet_names:
                if name not in wb.sheetnames:
                    continue

                ws = wb[name]
                data = []

                for row in ws.iter_rows():
                    row_data = []
                    for cell in row:
                        value = cell.value
                        if value is not None:
                            row_data.append(str(value))
                        else:
                            row_data.append("")
                    data.append(row_data)

                # Remove trailing empty rows
                while data and all(cell == "" for cell in data[-1]):
                    data.pop()

                sheets[name] = {
                    "data": data,
                    "rows": len(data),
                    "cols": max(len(row) for row in data) if data else 0,
                }

            return {
                "sheets": sheets,
                "sheet_names": wb.sheetnames,
                "active_sheet": wb.active.title if wb.active else None,
            }

        except Exception as e:
            logger.error(f"Error reading XLSX: {e}")
            return {"error": str(e)}

    def fill_xlsx_template(
        self,
        template_path: Union[str, Path, bytes],
        data: Dict[str, str],
        output_path: Optional[Union[str, Path]] = None,
    ) -> Dict[str, Any]:
        """
        Replace {{placeholders}} in an XLSX with data values.

        Args:
            template_path: Path to the template XLSX or bytes content.
            data: Dict of placeholder names to values.
            output_path: Where to save the result.

        Returns:
            Dict with success status.
        """
        try:
            from openpyxl import load_workbook
        except ImportError:
            return {"error": "openpyxl not installed. Run: pip install openpyxl"}

        try:
            if isinstance(template_path, bytes):
                wb = load_workbook(io.BytesIO(template_path))
            else:
                wb = load_workbook(template_path)

            for ws in wb.worksheets:
                for row in ws.iter_rows():
                    for cell in row:
                        if cell.value and isinstance(cell.value, str):
                            for key, value in data.items():
                                placeholder = "{{" + key + "}}"
                                if placeholder in cell.value:
                                    cell.value = cell.value.replace(placeholder, str(value))

            if output_path:
                wb.save(output_path)
                return {"success": True, "output_path": str(output_path)}
            else:
                buffer = io.BytesIO()
                wb.save(buffer)
                buffer.seek(0)
                return {"success": True, "content": buffer.read()}

        except Exception as e:
            logger.error(f"Error filling XLSX template: {e}")
            return {"error": str(e)}

    def xlsx_to_csv(
        self,
        file_path: Union[str, Path, bytes],
        sheet_name: Optional[str] = None,
        delimiter: str = ",",
    ) -> Dict[str, Any]:
        """
        Export XLSX sheet to CSV format.

        Args:
            file_path: Path to the XLSX file or bytes content.
            sheet_name: Sheet to export (uses first sheet if not specified).
            delimiter: CSV delimiter (default: comma).

        Returns:
            Dict containing the CSV text.
        """
        result = self.read_xlsx(file_path, sheet_name)
        if "error" in result:
            return result

        try:
            # Get the first sheet if not specified
            sheets = result.get("sheets", {})
            if not sheets:
                return {"error": "No sheets found in workbook"}

            target_sheet = sheet_name or list(sheets.keys())[0]
            if target_sheet not in sheets:
                return {"error": f"Sheet '{target_sheet}' not found"}

            data = sheets[target_sheet]["data"]

            # Convert to CSV
            lines = []
            for row in data:
                # Escape fields containing delimiter or quotes
                escaped_row = []
                for cell in row:
                    if delimiter in cell or '"' in cell or "\n" in cell:
                        cell = '"' + cell.replace('"', '""') + '"'
                    escaped_row.append(cell)
                lines.append(delimiter.join(escaped_row))

            return {
                "csv": "\n".join(lines),
                "sheet_name": target_sheet,
                "rows": len(data),
            }

        except Exception as e:
            logger.error(f"Error converting XLSX to CSV: {e}")
            return {"error": str(e)}

    # =========================================================================
    # PPTX Processing
    # =========================================================================

    def read_pptx(self, file_path: Union[str, Path, bytes]) -> Dict[str, Any]:
        """
        Extract slides, text, and speaker notes from a PPTX file.

        Args:
            file_path: Path to the PPTX file or bytes content.

        Returns:
            Dict containing slides data and metadata.
        """
        try:
            from pptx import Presentation
        except ImportError:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}

        try:
            if isinstance(file_path, bytes):
                prs = Presentation(io.BytesIO(file_path))
            else:
                prs = Presentation(file_path)

            slides = []
            for i, slide in enumerate(prs.slides, 1):
                slide_data = {
                    "number": i,
                    "title": "",
                    "content": [],
                    "notes": "",
                }

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        if shape.is_placeholder and shape.placeholder_format.type == 1:  # Title
                            slide_data["title"] = shape.text
                        else:
                            slide_data["content"].append(shape.text)

                # Extract speaker notes
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    if notes_slide.notes_text_frame:
                        slide_data["notes"] = notes_slide.notes_text_frame.text

                slides.append(slide_data)

            return {
                "slides": slides,
                "slide_count": len(slides),
                "width": prs.slide_width,
                "height": prs.slide_height,
            }

        except Exception as e:
            logger.error(f"Error reading PPTX: {e}")
            return {"error": str(e)}

    def fill_pptx_template(
        self,
        template_path: Union[str, Path, bytes],
        data: Dict[str, str],
        output_path: Optional[Union[str, Path]] = None,
    ) -> Dict[str, Any]:
        """
        Replace {{placeholders}} in a PPTX with data values.

        Args:
            template_path: Path to the template PPTX or bytes content.
            data: Dict of placeholder names to values.
            output_path: Where to save the result.

        Returns:
            Dict with success status.
        """
        try:
            from pptx import Presentation
        except ImportError:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}

        try:
            if isinstance(template_path, bytes):
                prs = Presentation(io.BytesIO(template_path))
            else:
                prs = Presentation(template_path)

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                for key, value in data.items():
                                    placeholder = "{{" + key + "}}"
                                    if placeholder in run.text:
                                        run.text = run.text.replace(placeholder, str(value))

                # Also check notes
                if slide.has_notes_slide:
                    notes = slide.notes_slide
                    if notes.notes_text_frame:
                        for para in notes.notes_text_frame.paragraphs:
                            for run in para.runs:
                                for key, value in data.items():
                                    placeholder = "{{" + key + "}}"
                                    if placeholder in run.text:
                                        run.text = run.text.replace(placeholder, str(value))

            if output_path:
                prs.save(output_path)
                return {"success": True, "output_path": str(output_path)}
            else:
                buffer = io.BytesIO()
                prs.save(buffer)
                buffer.seek(0)
                return {"success": True, "content": buffer.read()}

        except Exception as e:
            logger.error(f"Error filling PPTX template: {e}")
            return {"error": str(e)}

    def pptx_to_markdown(self, file_path: Union[str, Path, bytes]) -> Dict[str, Any]:
        """
        Extract PPTX content as a markdown outline.

        Args:
            file_path: Path to the PPTX file or bytes content.

        Returns:
            Dict containing the markdown text.
        """
        result = self.read_pptx(file_path)
        if "error" in result:
            return result

        try:
            lines = []

            for slide in result.get("slides", []):
                # Slide header
                lines.append(f"## Slide {slide['number']}")
                lines.append("")

                # Title
                if slide.get("title"):
                    lines.append(f"### {slide['title']}")
                    lines.append("")

                # Content
                for content in slide.get("content", []):
                    if content.strip():
                        # Check if it looks like a bullet point
                        for line in content.split("\n"):
                            if line.strip():
                                lines.append(f"- {line.strip()}")
                        lines.append("")

                # Notes
                if slide.get("notes"):
                    lines.append("**Speaker Notes:**")
                    lines.append(f"> {slide['notes']}")
                    lines.append("")

                lines.append("---")
                lines.append("")

            return {
                "markdown": "\n".join(lines),
                "slide_count": result.get("slide_count", 0),
            }

        except Exception as e:
            logger.error(f"Error converting PPTX to markdown: {e}")
            return {"error": str(e)}


# Singleton instance
_processor: Optional[OfficeProcessor] = None


def get_office_processor() -> OfficeProcessor:
    """Get the singleton OfficeProcessor instance."""
    global _processor
    if _processor is None:
        _processor = OfficeProcessor()
    return _processor
